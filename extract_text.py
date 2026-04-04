#!/usr/bin/env python3
"""Extract text from a local file for the Webex MCP server."""

from __future__ import annotations

import json
import mimetypes
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".html",
    ".htm",
    ".js",
    ".ts",
    ".mjs",
    ".cjs",
    ".py",
    ".log",
    ".yaml",
    ".yml",
    ".ini",
    ".cfg",
  }


def read_text_file(path: Path) -> dict[str, object]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return {"kind": "text", "text": text}


def extract_pdf(path: Path) -> dict[str, object]:
    reader = PdfReader(str(path))
    pages: list[str] = []

    for page in reader.pages:
        pages.append(page.extract_text() or "")

    return {"kind": "pdf", "text": "\n\n".join(pages).strip()}


def extract_docx(path: Path) -> dict[str, object]:
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("word/document.xml")

    root = ET.fromstring(xml)
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    chunks = [node.text for node in root.findall(".//w:t", namespace) if node.text]

    return {"kind": "docx", "text": "\n".join(chunks).strip()}


def extract_pptx(path: Path) -> dict[str, object]:
    presentation = Presentation(str(path))
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
        if parts:
            slides.append(f"[Slide {index}]\n" + "\n".join(parts))

    return {"kind": "pptx", "text": "\n\n".join(slides).strip()}


def extract_xlsx(path: Path) -> dict[str, object]:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    rendered: list[str] = []

    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if cell is None else str(cell) for cell in row]
            if any(cells):
                rows.append("\t".join(cells).rstrip())
        if rows:
            rendered.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))

    return {"kind": "xlsx", "text": "\n\n".join(rendered).strip()}


def extract(path: Path, mime_type: str | None = None) -> dict[str, object]:
    suffix = path.suffix.lower()
    detected_mime = mime_type or mimetypes.guess_type(path.name)[0] or "application/octet-stream"

    if suffix in TEXT_EXTENSIONS or detected_mime.startswith("text/") or detected_mime in {
        "application/json",
        "application/xml",
        "application/javascript",
      }:
        return extract_text_payload(read_text_file(path), detected_mime)

    if suffix == ".pdf" or detected_mime == "application/pdf":
        return extract_text_payload(extract_pdf(path), detected_mime)

    if suffix == ".docx" or detected_mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_payload(extract_docx(path), detected_mime)

    if suffix == ".pptx" or detected_mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_payload(extract_pptx(path), detected_mime)

    if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"} or detected_mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return extract_text_payload(extract_xlsx(path), detected_mime)

    raise RuntimeError(f"Unsupported extraction format for {path.name} ({detected_mime})")


def extract_text_payload(result: dict[str, object], mime_type: str) -> dict[str, object]:
    text = str(result.get("text") or "").strip()
    return {
        "mimeType": mime_type,
        "kind": result.get("kind") or "text",
        "text": text,
        "textLength": len(text),
    }


def main() -> None:
    if len(sys.argv) < 2:
        raise SystemExit("usage: extract_text.py <path> [mimeType]")

    path = Path(sys.argv[1]).expanduser().resolve()
    mime_type = sys.argv[2] if len(sys.argv) > 2 else None
    result = extract(path, mime_type)
    print(json.dumps(result))


if __name__ == "__main__":
    main()
